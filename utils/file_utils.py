# utils/file_utils.py
"""
File reading utilities for NarrativeNexus AI.
Supports: PDF, DOCX, PPTX, EPUB, CSV, Excel, JSON, HTML, RTF, XML, TXT, Images (OCR)
Includes file size validation, sanitization, and graceful error handling.
"""
import pandas as pd
import json
from bs4 import BeautifulSoup
from docx import Document
import pdfplumber
from pptx import Presentation
import striprtf
from lxml import etree as ET
from PIL import Image
import ebooklib
from ebooklib import epub
import tempfile
import os
from pathlib import Path
from utils.config import MAX_FILE_SIZE_MB, FEATURE_OCR


# ── File Size Validation ──────────────────────────────────────────────────────
def validate_file(file) -> tuple[bool, str]:
    """Validate file before processing. Returns (is_valid, error_message)."""
    if file is None:
        return False, "No file provided."
    
    # Check file size
    file.seek(0, 2)  # Seek to end
    size_bytes = file.tell()
    file.seek(0)  # Reset
    size_mb = size_bytes / (1024 * 1024)
    
    if size_mb > MAX_FILE_SIZE_MB:
        return False, f"File too large ({size_mb:.1f} MB). Maximum is {MAX_FILE_SIZE_MB} MB."
    
    return True, ""


# ── Sanitize Filename ─────────────────────────────────────────────────────────
def sanitize_filename(name: str) -> str:
    """Sanitize filename to prevent path traversal attacks."""
    name = Path(name).name  # Strip any directory components
    # Remove dangerous characters
    safe_chars = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789._- "
    return "".join(c for c in name if c in safe_chars)[:255]


# ── OCR Handler (Images) ──────────────────────────────────────────────────────
def read_image(file) -> str:
    """Extract text from image using Tesseract OCR."""
    if not FEATURE_OCR:
        return "[OCR disabled]"
    
    try:
        import pytesseract
        image = Image.open(file)
        text = pytesseract.image_to_string(image)
        return text.strip() or "[Image contains no readable text]"
    except ImportError:
        return "[OCR unavailable: pytesseract not installed]"
    except Exception as e:
        # Try without pytesseract — just return a placeholder
        return f"[Image OCR Error: {e}. Ensure Tesseract is installed.]"


# ── PDF Handler ───────────────────────────────────────────────────────────────
def read_pdf(file) -> str:
    """Extract text from PDF using pdfplumber."""
    try:
        text_parts = []
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                extracted = page.extract_text() or ""
                if extracted.strip():
                    text_parts.append(extracted)
        
        result = "\n".join(text_parts)
        
        # If PDF has no extractable text, it might be scanned — try OCR
        if not result.strip() and FEATURE_OCR:
            try:
                import pytesseract
                with pdfplumber.open(file) as pdf:
                    ocr_parts = []
                    for page in pdf.pages:
                        img = page.to_image(resolution=150).original
                        ocr_text = pytesseract.image_to_string(img)
                        if ocr_text.strip():
                            ocr_parts.append(ocr_text)
                    return "\n".join(ocr_parts) or "[PDF: no readable text found]"
            except Exception:
                return "[PDF appears to be scanned/image-based. OCR not available.]"
        
        return result or "[PDF: no readable text found]"
    except Exception as e:
        return f"[PDF read error: {e}]"


# ── DOCX Handler ──────────────────────────────────────────────────────────────
def read_docx(file) -> str:
    """Extract text from DOCX file including tables."""
    try:
        doc = Document(file)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Also extract table content
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"[DOCX read error: {e}]"


# ── PPTX Handler ──────────────────────────────────────────────────────────────
def read_pptx(file) -> str:
    """Extract text from PowerPoint presentation."""
    try:
        text = []
        prs = Presentation(file)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text.append(f"[Slide {i+1}]\n" + "\n".join(slide_text))
        return "\n\n".join(text)
    except Exception as e:
        return f"[PPTX read error: {e}]"


# ── EPUB Handler ──────────────────────────────────────────────────────────────
def read_epub(file) -> str:
    """Extract text from EPUB ebook."""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".epub") as tmp:
            tmp.write(file.getvalue())
            tmp_path = tmp.name
        
        book = epub.read_epub(tmp_path)
        chapters = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                text = soup.get_text(separator="\n")
                if text.strip():
                    chapters.append(text)
        
        os.remove(tmp_path)
        return "\n\n".join(chapters) or "[EPUB: no readable content]"
    except Exception as e:
        return f"[EPUB read error: {e}]"


# ── Text/CSV/JSON/HTML ────────────────────────────────────────────────────────
def read_txt(file) -> str:
    """Read plain text file."""
    try:
        return file.read().decode("utf-8", errors="ignore").strip()
    except Exception:
        return "[Text decode error]"


def read_csv(file) -> str:
    """Read CSV and convert to text representation."""
    try:
        df = pd.read_csv(file)
        # Include headers as part of the text
        header_text = "Columns: " + ", ".join(df.columns.tolist())
        row_text = "\n".join(df.astype(str).agg(" | ".join, axis=1))
        return f"{header_text}\n\n{row_text}"
    except Exception as e:
        return f"[CSV error: {e}]"


def read_excel(file) -> str:
    """Read Excel file and convert to text."""
    try:
        # Try to read all sheets
        xl = pd.ExcelFile(file)
        all_text = []
        for sheet in xl.sheet_names[:5]:  # Limit to 5 sheets
            df = pd.read_excel(file, sheet_name=sheet)
            all_text.append(f"[Sheet: {sheet}]")
            all_text.append("\n".join(df.astype(str).agg(" | ".join, axis=1)))
        return "\n\n".join(all_text)
    except Exception as e:
        return f"[Excel error: {e}]"


def read_json(file) -> str:
    """Read JSON and convert to readable text."""
    try:
        data = json.load(file)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception as e:
        return f"[JSON error: {e}]"


def read_html(file) -> str:
    """Extract text from HTML."""
    try:
        soup = BeautifulSoup(file.read(), "html.parser")
        # Remove script/style tags
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except Exception as e:
        return f"[HTML error: {e}]"


def read_rtf(file) -> str:
    """Read RTF file."""
    try:
        return striprtf.rtf_to_text(file.read().decode("utf-8", errors="ignore"))
    except Exception as e:
        return f"[RTF error: {e}]"


def read_xml(file) -> str:
    """Read XML file."""
    try:
        tree = ET.parse(file)
        root = tree.getroot()
        return ET.tostring(root, encoding="unicode", method="text")
    except Exception as e:
        return f"[XML error: {e}]"


def read_markdown(file) -> str:
    """Read Markdown file."""
    try:
        content = file.read().decode("utf-8", errors="ignore")
        # Strip markdown syntax for NLP processing
        soup = BeautifulSoup(content, "html.parser")
        return soup.get_text(separator="\n", strip=True) or content
    except Exception as e:
        return f"[Markdown error: {e}]"


# ── Master Handler ─────────────────────────────────────────────────────────────
def read_file(file) -> str:
    """
    Master file dispatcher. Routes to appropriate handler based on file extension.
    Returns extracted text content.
    """
    name = file.name.lower()
    
    if name.endswith(".pdf"):
        return read_pdf(file)
    elif name.endswith((".docx", ".doc")):
        return read_docx(file)
    elif name.endswith((".pptx", ".ppt")):
        return read_pptx(file)
    elif name.endswith(".epub"):
        return read_epub(file)
    elif name.endswith(".txt"):
        return read_txt(file)
    elif name.endswith(".csv"):
        return read_csv(file)
    elif name.endswith((".xlsx", ".xls")):
        return read_excel(file)
    elif name.endswith(".json"):
        return read_json(file)
    elif name.endswith((".html", ".htm")):
        return read_html(file)
    elif name.endswith(".rtf"):
        return read_rtf(file)
    elif name.endswith(".xml"):
        return read_xml(file)
    elif name.endswith(".md"):
        return read_markdown(file)
    elif name.endswith((".png", ".jpg", ".jpeg", ".tiff", ".bmp")):
        return read_image(file)
    else:
        return f"[Unsupported format: {name.split('.')[-1]}]"
